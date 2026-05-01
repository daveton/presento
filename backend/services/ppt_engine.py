from typing import Dict, Any, List
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os
import uuid

# 品牌色
PRIMARY_COLOR = RgbColor(0x4D, 0x77, 0xFF)  # #4D77FF
TEXT_COLOR = RgbColor(0x1E, 0x29, 0x3B)     # #1E293B
SUBTEXT_COLOR = RgbColor(0x47, 0x55, 0x69)  # #475569

async def create_ppt_file(content: Dict[str, Any]) -> str:
    """
    创建 PPT 文件并返回下载链接
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slides_data = content.get("slides", [])
    
    for slide_data in slides_data:
        slide_type = slide_data.get("type", "content")
        
        if slide_type == "cover":
            _create_cover_slide(prs, slide_data)
        else:
            _create_content_slide(prs, slide_data)
    
    # 确保输出目录存在
    output_dir = "/tmp/presento"
    os.makedirs(output_dir, exist_ok=True)
    
    # 生成唯一文件名
    filename = f"{uuid.uuid4().hex}.pptx"
    filepath = os.path.join(output_dir, filename)
    
    prs.save(filepath)
    
    # 返回相对路径，实际部署时转换为完整 URL
    return f"/download/{filename}"


def _create_cover_slide(prs: Presentation, data: Dict[str, Any]):
    """创建封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景色块
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(1), Inches(2), Inches(11.333), Inches(3.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(0xF0, 0xF4, 0xFF)  # #F0F4FF
    shape.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.5), Inches(10.333), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = _truncate_text(data.get("title", ""), 22)
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if data.get("subtitle"):
        subtitle_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(4), Inches(10.333), Inches(0.8)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = _truncate_text(data["subtitle"], 40)
        p.font.size = Pt(20)
        p.font.color.rgb = SUBTEXT_COLOR
        p.alignment = PP_ALIGN.CENTER


def _create_content_slide(prs: Presentation, data: Dict[str, Any]):
    """创建内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 页面标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.8), Inches(11.333), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = _truncate_text(data.get("title", ""), 18)
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_COLOR
    
    # 要点列表
    points = data.get("points", [])
    if points:
        # 根据数量确定字号
        font_size = _get_font_size(len(points))
        
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(11.333), Inches(5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(points[:5]):  # 最多5条
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = f"• {_truncate_text(point, 20)}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = SUBTEXT_COLOR
            p.space_after = Pt(16)
            p.level = 0


def _get_font_size(point_count: int) -> int:
    """根据要点数量确定字号"""
    if point_count <= 3:
        return 28
    elif point_count <= 4:
        return 24
    else:
        return 20


def _truncate_text(text: str, max_chars: int) -> str:
    """截断文本"""
    if len(text) > max_chars:
        return text[:max_chars-1] + "…"
    return text
