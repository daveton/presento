from typing import Dict, Any, List
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os
import uuid

# ==============================
# 配置
# ==============================

OUTPUT_DIR = "/tmp/presento"
MAX_TITLE_LEN = 22
MAX_POINTS = 5
MAX_POINT_LEN = 20

FONT_SIZES = {
    3: 28,
    5: 24,
    8: 20
}

# 品牌色
PRIMARY_COLOR = RgbColor(0x4D, 0x77, 0xFF)  # #4D77FF
TEXT_COLOR = RgbColor(0x1E, 0x29, 0x3B)     # #1E293B
SUBTEXT_COLOR = RgbColor(0x47, 0x55, 0x69)  # #475569

# ==============================
# 工具函数
# ==============================

def trim_text(text, max_len):
    return text[:max_len]

def clean_text(text):
    return text.replace("，", "").replace("。", "")

def split_points(points):
    """自动分页"""
    result = []
    for i in range(0, len(points), MAX_POINTS):
        result.append(points[i:i + MAX_POINTS])
    return result

def get_font_size(n):
    if n <= 3:
        return FONT_SIZES[3]
    elif n <= 5:
        return FONT_SIZES[5]
    else:
        return FONT_SIZES[8]

def rewrite_point(text):
    """
    简单内容优化
    """
    text = clean_text(text)

    # 强制短句
    if len(text) > MAX_POINT_LEN:
        text = text[:MAX_POINT_LEN]

    # 强化节奏（简单规则）
    if "因为" in text:
        text = text.replace("因为", "原因：")

    return text


# ==============================
# 核心规则引擎
# ==============================

def enforce_rules(slides):
    """强制规则引擎 - 确保所有输出符合规范"""
    new_slides = []

    for slide in slides:
        slide_type = slide.get("type", "content")

        # 封面
        if slide_type == "cover":
            new_slides.append({
                "type": "cover",
                "title": trim_text(slide.get("title", ""), MAX_TITLE_LEN),
                "subtitle": trim_text(slide.get("subtitle", ""), 40)
            })
            continue

        # 内容页
        title = trim_text(slide.get("title", ""), MAX_TITLE_LEN)
        points = slide.get("points", [])

        # 重写 + 裁剪
        points = [rewrite_point(p) for p in points]
        points = [p[:MAX_POINT_LEN] for p in points]

        # 自动分页
        split = split_points(points)

        for part in split:
            new_slides.append({
                "type": "content",
                "title": title,
                "points": part
            })

    return new_slides[:10]  # 限制总页数


# ==============================
# 质量评分
# ==============================

def score_slides(slides):
    """评分系统 - 满分100，合格线60"""
    score = 100

    for slide in slides:
        if slide["type"] == "content":
            pts = slide.get("points", [])

            if len(pts) < 2:
                score -= 20

            for p in pts:
                if len(p) > MAX_POINT_LEN:
                    score -= 5

    return score


# ==============================
# 主函数
# ==============================

async def create_ppt_file(content: Dict[str, Any]) -> str:
    """
    创建 PPT 文件并返回下载链接
    完整流程：规则引擎 -> 质量评分 -> 渲染
    """
    slides_data = content.get("slides", [])
    
    # 1️⃣ 规则引擎
    slides_data = enforce_rules(slides_data)
    
    # 2️⃣ 质量评分
    quality_score = score_slides(slides_data)
    if quality_score < 60:
        print(f"Warning: Quality score {quality_score} is below threshold")
        # 继续渲染，但记录警告
    
    # 3️⃣ 创建PPT
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for slide_data in slides_data:
        slide_type = slide_data.get("type", "content")
        
        if slide_type == "cover":
            _create_cover_slide(prs, slide_data)
        else:
            _create_content_slide(prs, slide_data)
    
    # 确保输出目录存在
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    
    # 生成唯一文件名
    filename = f"{uuid.uuid4().hex}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    
    prs.save(filepath)
    
    return f"/download/{filename}"


def _create_cover_slide(prs: Presentation, data: Dict[str, Any]):
    """创建封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景色块
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(1), Inches(2), Inches(11.333), Inches(3.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(0xF0, 0xF4, 0xFF)  # #F0F4FF
    shape.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.5), Inches(10.333), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = _truncate_text(data.get("title", ""), 22)
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if data.get("subtitle"):
        subtitle_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(4), Inches(10.333), Inches(0.8)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = _truncate_text(data["subtitle"], 40)
        p.font.size = Pt(20)
        p.font.color.rgb = SUBTEXT_COLOR
        p.alignment = PP_ALIGN.CENTER


def _create_content_slide(prs: Presentation, data: Dict[str, Any]):
    """创建内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 页面标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.8), Inches(11.333), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = _truncate_text(data.get("title", ""), 18)
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_COLOR
    
    # 要点列表
    points = data.get("points", [])
    if points:
        # 根据数量确定字号
        font_size = _get_font_size(len(points))
        
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(11.333), Inches(5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(points[:5]):  # 最多5条
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = f"• {_truncate_text(point, 20)}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = SUBTEXT_COLOR
            p.space_after = Pt(16)
            p.level = 0


def _get_font_size(point_count: int) -> int:
    """根据要点数量确定字号"""
    if point_count <= 3:
        return 28
    elif point_count <= 4:
        return 24
    else:
        return 20


def _truncate_text(text: str, max_chars: int) -> str:
    """截断文本"""
    if len(text) > max_chars:
        return text[:max_chars-1] + "…"
    return text
