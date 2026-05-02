"""
模板渲染器 - Lite 版
加载模板作为base，清空内容，填充新内容
"""

import os
import uuid
from pptx import Presentation
from pptx.util import Pt
from typing import Dict, Any, List

from .template_config import get_template_config, choose_layout


def create_ppt_with_template(content: Dict[str, Any], template_id: str = "business") -> str:
    """
    基于模板生成PPT
    流程：加载模板 → 清空slides → 根据内容创建新slides
    """
    # 获取模板配置
    config = get_template_config(template_id)
    template_file = config["file"]

    # 检查模板文件是否存在
    if not os.path.exists(template_file):
        print(f"[Template] Template not found: {template_file}, using default renderer")
        from .renderer import create_ppt_file
        return create_ppt_file(content)

    # 加载模板
    prs = Presentation(template_file)

    # 清空原有slides（保留母版）
    # 注意：不能删除所有slide，至少需要保留一个来复制母版
    # 这里采用：保留母版，删除所有内容slide
    slide_layouts = prs.slide_layouts

    # 收集要删除的slides（保留第一个用于复制母版信息）
    slides_to_remove = list(prs.slides)[1:] if len(prs.slides) > 1 else []

    # 删除多余的slides
    for slide in slides_to_remove:
        rId = slide.part.rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[prs.slides.index(slide)]

    # 如果没有slides，创建一个空白slide占位
    if len(prs.slides) == 0:
        blank_layout = slide_layouts[6] if len(slide_layouts) > 6 else slide_layouts[-1]
        prs.slides.add_slide(blank_layout)

    # 获取slides数据
    slides_data = content.get("slides", [])

    # 为每个内容创建slide
    for i, slide_data in enumerate(slides_data):
        # 选择layout
        layout_name = choose_layout(slide_data, template_id)
        layout_index = config["layouts"].get(layout_name, 1)

        # 确保layout index有效
        if layout_index >= len(slide_layouts):
            layout_index = 1 if len(slide_layouts) > 1 else 0

        slide_layout = slide_layouts[layout_index]

        # 添加新slide
        slide = prs.slides.add_slide(slide_layout)

        # 填充内容
        fill_slide_content(slide, slide_data)

    # 删除第一个占位slide（如果创建了多个）
    if len(prs.slides) > len(slides_data):
        first_slide = prs.slides[0]
        rId = first_slide.part.rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # 保存文件
    output_dir = "/tmp/presento"
    os.makedirs(output_dir, exist_ok=True)

    filename = f"presento_{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join(output_dir, filename)

    prs.save(filepath)

    download_url = f"/api/download/{filename}"
    print(f"[Template] Created PPT with template '{template_id}': {download_url}")

    return download_url


def fill_slide_content(slide, slide_data: Dict[str, Any]):
    """
    填充slide内容到占位符
    """
    title = slide_data.get("title", "")
    points = slide_data.get("points", [])

    # 尝试设置标题
    if slide.shapes.title:
        slide.shapes.title.text = title[:22]  # 限制长度

    # 查找内容占位符（通常是body/placeholder）
    for shape in slide.shapes.placeholders:
        # 1 = title, 2 = body/content
        if shape.placeholder_format.idx == 1:
            shape.text = title[:22]
        elif shape.placeholder_format.idx == 2 and points:
            # 填充要点
            text_frame = shape.text_frame
            text_frame.clear()

            for i, point in enumerate(points[:5]):  # 最多5条
                p = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
                p.text = point[:20]  # 限制20字
                p.level = 0

    # 如果没有找到标准占位符，尝试其他形状
    if not points:
        return

    # 查找非标题的文本框
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            text_frame = shape.text_frame
            text_frame.clear()

            for i, point in enumerate(points[:5]):
                p = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
                p.text = point[:20]

            break  # 只填充第一个合适的文本框
